import os
import re
import docx
import pdfplumber
import pandas as pd
from openai import OpenAI
from pptx import Presentation
import chardet
from datetime import datetime

# 配置OpenAI客户端
client = OpenAI(
    base_url='https://api-inference.modelscope.cn/v1/',
    api_key='your_api_key',
)

# 额外配置
extra_body = {
    "enable_thinking": True,
}

def detect_language(text):
    """简单的语言检测"""
    if not text.strip():
        return 'en'
    
    # 统计中文字符数量
    chinese_chars = sum(1 for char in text if '\u4e00' <= char <= '\u9fff')
    # 统计英文字符数量
    english_chars = sum(1 for char in text if char.isalpha() and char.isascii())
    
    # 如果中文字符占比超过30%，判定为中文
    if chinese_chars / max(len(text), 1) > 0.3:
        return 'zh-cn'
    
    # 如果英文字符占比超过50%，判定为英文
    if english_chars / max(len(text), 1) > 0.5:
        return 'en'
    
    return 'en'

def extract_text_from_pdf(file_path):
    """从PDF提取文本"""
    with pdfplumber.open(file_path) as pdf:
        text = ""
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_docx(file_path):
    """从Word文档提取文本"""
    doc = docx.Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

def extract_text_from_excel(file_path):
    """从Excel提取文本"""
    xls = pd.ExcelFile(file_path)
    text = ""
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name)
        text += f"工作表: {sheet_name}\n"
        text += df.to_string(na_rep='nan', index=False) + "\n\n"
    return text

def extract_text_from_pptx(file_path):
    """从PowerPoint提取文本"""
    prs = Presentation(file_path)
    text = ""
    for i, slide in enumerate(prs.slides, 1):
        text += f"幻灯片 {i}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        text += "\n"
    return text

def extract_text_from_srt(file_path):
    """从SRT字幕文件提取文本"""
    try:
        # 尝试检测文件编码
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            encoding = chardet.detect(raw_data)['encoding']
        
        # 使用检测到的编码读取文件
        with open(file_path, 'r', encoding=encoding or 'utf-8') as f:
            content = f.read()
        
        # 解析SRT格式
        text = ""
        lines = content.strip().split('\n')
        i = 0
        while i < len(lines):
            # 跳过序号行
            if lines[i].strip().isdigit():
                i += 1
                # 跳过时间戳行
                if i < len(lines) and '-->' in lines[i]:
                    i += 1
                    # 读取字幕文本
                    subtitle_text = ""
                    while i < len(lines) and lines[i].strip() != "":
                        subtitle_text += lines[i] + " "
                        i += 1
                    text += subtitle_text.strip() + "\n"
            i += 1
        return text
    except Exception as e:
        print(f"读取SRT文件失败: {e}")
        return ""

def extract_text_from_txt(file_path):
    """从TXT文本文件提取文本"""
    try:
        # 尝试检测文件编码
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            encoding = chardet.detect(raw_data)['encoding']
        
        # 使用检测到的编码读取文件
        with open(file_path, 'r', encoding=encoding or 'utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"读取TXT文件失败: {e}")
        return ""

def get_file_extractor(file_path):
    """根据文件扩展名获取对应的文本提取函数"""
    file_ext = file_path.lower().split('.')[-1]
    extractors = {
        'pdf': extract_text_from_pdf,
        'docx': extract_text_from_docx,
        'doc': extract_text_from_docx,
        'xlsx': extract_text_from_excel,
        'xls': extract_text_from_excel,
        'pptx': extract_text_from_pptx,
        'ppt': extract_text_from_pptx,
        'srt': extract_text_from_srt,
        'txt': extract_text_from_txt,
    }
    return extractors.get(file_ext)

def summarize_and_interpret(text, language, file_name):
    """使用Qwen3总结和解读文本"""
    if language == 'zh-cn':
        system_prompt = f"""你是专业的文档解读专家。请详细总结并解读文件'{file_name}'的内容，要求：
        1. 提炼关键要点和重要数据；
        2. 分析潜在含义和实际应用建议；
        3. 用结构化格式呈现，确保可读性。"""
    else:
        system_prompt = f"""You are a professional document interpreter. Please provide a detailed Chinese summary and interpretation of the content from file '{file_name}', including:
        1. Key points and important data;
        2. Potential implications and practical suggestions;
        3. Structured formatting for readability."""
    
    # 限制输入长度
    truncated_text = text[:100000]
    response = client.chat.completions.create(
        model='Qwen/Qwen3-235B-A22B',
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": truncated_text}
        ],
        stream=True,
        extra_body=extra_body
    )
    
    summary = ""
    done_thinking = False
    for chunk in response:
        answer_chunk = chunk.choices[0].delta.content
        if answer_chunk:
            if not done_thinking:
                summary += "\n\n=== 文档解读结果 ===\n"
                done_thinking = True
            summary += answer_chunk
    return summary.strip()

def parse_inline_formatting(text):
    """解析内联格式（粗体、斜体等）"""
    parts = []
    current_pos = 0
    
    # 查找粗体标记 **text**
    while current_pos < len(text):
        bold_start = text.find('**', current_pos)
        if bold_start == -1:
            if current_pos < len(text):
                parts.append(('normal', text[current_pos:]))
            break
        
        # 添加粗体标记前的普通文本
        if bold_start > current_pos:
            parts.append(('normal', text[current_pos:bold_start]))
        
        # 查找粗体结束标记
        bold_end = text.find('**', bold_start + 2)
        if bold_end == -1:
            parts.append(('normal', text[bold_start:]))
            break
        
        # 添加粗体文本
        bold_text = text[bold_start + 2:bold_end]
        if bold_text:
            parts.append(('bold', bold_text))
        
        current_pos = bold_end + 2
    
    return parts

def add_formatted_paragraph(doc, text, style=None):
    """添加带格式的段落"""
    para = doc.add_paragraph()
    if style:
        para.style = style
    
    parts = parse_inline_formatting(text)
    for format_type, content in parts:
        run = para.add_run(content)
        if format_type == 'bold':
            run.bold = True
    
    return para

def save_to_word(content, output_file_path):
    """保存内容到Word文档（解析Markdown格式并转换为格式化的Word文档）"""
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    
    doc = docx.Document()
    
    # 设置文档样式
    style = doc.styles['Normal']
    style.font.name = '微软雅黑'
    style.font.size = Pt(11)
    
    # 按行分割内容
    lines = content.split('\n')
    current_table = None
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # 处理标题
        if line.startswith('### '):
            heading = doc.add_heading(line[4:], level=3)
            heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
            current_table = None
        elif line.startswith('## '):
            heading = doc.add_heading(line[3:], level=2)
            heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
            current_table = None
        elif line.startswith('# '):
            heading = doc.add_heading(line[2:], level=1)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            current_table = None
        elif line.startswith('===') and '===' in line:
            title_text = line.replace('=', '').strip()
            if title_text:
                heading = doc.add_heading(title_text, level=1)
                heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            current_table = None
        elif line.startswith('- '):
            list_text = line[2:]
            add_formatted_paragraph(doc, list_text, 'List Bullet')
            current_table = None
        elif '|' in line and line.count('|') >= 3:
            # 表格行
            if '---' not in line:
                cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                if cells:
                    if current_table is None:
                        current_table = doc.add_table(rows=1, cols=len(cells))
                        current_table.style = 'Table Grid'
                        # 添加表头
                        for i, cell_text in enumerate(cells):
                            cell = current_table.cell(0, i)
                            cell_text = cell_text.replace('<br>', '\n')
                            cell.text = cell_text
                            # 设置表头为粗体
                            for paragraph in cell.paragraphs:
                                for run in paragraph.runs:
                                    run.bold = True
                    else:
                        # 添加数据行
                        row = current_table.add_row()
                        for i, cell_text in enumerate(cells):
                            if i < len(row.cells):
                                cell_text = cell_text.replace('<br>', '\n')
                                cell = row.cells[i]
                                cell.text = ''
                                
                                # 按换行符分割并添加段落
                                cell_lines = cell_text.split('\n')
                                for j, cell_line in enumerate(cell_lines):
                                    if j == 0:
                                        para = cell.paragraphs[0]
                                    else:
                                        para = cell.add_paragraph()
                                    
                                    # 解析内联格式
                                    parts = parse_inline_formatting(cell_line)
                                    for format_type, content in parts:
                                        run = para.add_run(content)
                                        if format_type == 'bold':
                                            run.bold = True
        elif line.startswith('---'):
            current_table = None
            doc.add_paragraph('_' * 50)
        elif line.startswith('::'):
            continue
        else:
            if line:
                add_formatted_paragraph(doc, line)
                current_table = None
    
    doc.save(output_file_path)

def clean_filename(filename):
    """清理文件名，移除特殊字符但保留中文字符"""
    base_name = os.path.splitext(filename)[0]
    cleaned = re.sub(r'[《》<>"|*?\\/:]+', '', base_name)
    cleaned = re.sub(r'\s+', '', cleaned)
    cleaned = re.sub(r'_+', '_', cleaned)
    cleaned = cleaned.strip('_')
    return cleaned

def process_single_document(file_path, output_folder):
    """处理单个文档文件"""
    try:
        print(f"开始处理文件: {file_path}")
        
        # 获取文件提取函数
        extractor = get_file_extractor(file_path)
        if not extractor:
            raise Exception(f"不支持的文件格式: {file_path}")
        
        # 提取文本
        text = extractor(file_path)
        if not text or not text.strip():
            raise Exception("文件内容为空或无法提取文本")
        
        print(f"成功提取文本，长度: {len(text)}")
        
        # 检测语言
        language = detect_language(text)
        print(f"检测到语言: {language}")
        
        # 生成总结
        file_name = os.path.basename(file_path)
        summary = summarize_and_interpret(text, language, file_name)
        
        # 保存结果
        cleaned_name = clean_filename(file_name)
        current_date = datetime.now().strftime("%Y%m%d")
        output_file_name = f"{cleaned_name}_{current_date}.docx"
        output_file_path = os.path.join(output_folder, output_file_name)
        save_to_word(summary, output_file_path)
        
        print(f"处理完成，结果保存到: {output_file_path}")
        return True, "处理成功"
        
    except Exception as e:
        error_msg = f"处理文件失败: {str(e)}"
        print(error_msg)
        return False, error_msg

def process_single_document_with_name(file_path, output_folder, original_filename=None):
    """处理单个文档文件，可以指定用于命名的原始文件名"""
    try:
        print(f"开始处理文件: {file_path}")
        
        # 获取文件提取函数
        extractor = get_file_extractor(file_path)
        if not extractor:
            raise Exception(f"不支持的文件格式: {file_path}")
        
        # 提取文本
        text = extractor(file_path)
        if not text or not text.strip():
            raise Exception("文件内容为空或无法提取文本")
        
        print(f"成功提取文本，长度: {len(text)}")
        
        # 检测语言
        language = detect_language(text)
        print(f"检测到语言: {language}")
        
        # 生成总结
        display_name = original_filename if original_filename else os.path.basename(file_path)
        summary = summarize_and_interpret(text, language, display_name)
        
        # 保存结果
        name_for_output = original_filename if original_filename else os.path.basename(file_path)
        cleaned_name = clean_filename(name_for_output)
        current_date = datetime.now().strftime("%Y%m%d")
        output_file_name = f"{cleaned_name}_{current_date}.docx"
        output_file_path = os.path.join(output_folder, output_file_name)
        save_to_word(summary, output_file_path)
        
        print(f"处理完成，结果保存到: {output_file_path}")
        return True, "处理成功"
        
    except Exception as e:
        error_msg = f"处理文件失败: {str(e)}"
        print(error_msg)
        return False, error_msg 